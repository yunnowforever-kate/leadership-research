"""PPT 작성 에이전트: python-pptx + Brandlogy 디자인 시스템 (16:9)"""
import os
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

BASE_DIR = os.path.dirname(os.path.dirname(__file__))
OUTPUT_DIR = os.environ.get("RESEARCH_OUTPUT_DIR", os.path.join(BASE_DIR, "outputs"))

# ── 슬라이드 치수 (EMU) ───────────────────────────────────────────────────────
W  = Inches(13.333)
H  = Inches(7.5)
MH = Inches(0.5)   # 좌우 마진

# ── 존 좌표 ──────────────────────────────────────────────────────────────────
HEADER_TOP    = Inches(0.4)
HEADER_H      = Inches(0.3)
HEADLINE_TOP  = Inches(1.0)
HEADLINE_H    = Inches(0.75)
SUBTITLE_TOP  = Inches(1.63)
SUBTITLE_H    = Inches(0.4)
BODY_TOP      = Inches(2.39)
BODY_H        = Inches(4.46)   # 6.85 - 2.39
FOOTER_TOP    = Inches(7.05)
FOOTER_H      = Inches(0.25)
CONTENT_W     = Inches(12.333)

# ── 컬러 ─────────────────────────────────────────────────────────────────────
C = {
    "headline":  RGBColor(0x22, 0x22, 0x22),
    "subtitle":  RGBColor(0x45, 0x51, 0x5e),
    "body":      RGBColor(0x22, 0x22, 0x22),
    "muted":     RGBColor(0x8e, 0x8e, 0x93),
    "blue":      RGBColor(0x14, 0x56, 0xf0),
    "blue2":     RGBColor(0x3b, 0x82, 0xf6),
    "blue_lt":   RGBColor(0x60, 0xa5, 0xfa),
    "surface":   RGBColor(0xf2, 0xf3, 0xf5),
    "border":    RGBColor(0xe5, 0xe7, 0xeb),
    "white":     RGBColor(0xff, 0xff, 0xff),
    "dark":      RGBColor(0x18, 0x1e, 0x25),
}

FONT = "Pretendard"


# ── 헬퍼 ─────────────────────────────────────────────────────────────────────

def _rgb_hex(r: RGBColor) -> str:
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"


def _add_textbox(slide, left, top, width, height,
                 text="", size=14, bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True) -> object:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C["body"]
    return txBox


def _add_rect(slide, left, top, width, height,
              fill_color=None, line_color=None, line_width=None,
              radius_emu=None) -> object:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    if radius_emu:
        sp = shape._element
        spPr = sp.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            spPr.remove(prstGeom)
        prstGeom = etree.SubElement(spPr, qn("a:prstGeom"), attrib={"prst": "roundRect"})
        avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        # radius ratio: radius_emu / min(width, height) * 2 * 100000 → capped at 50000
        ratio = min(int(radius_emu / min(width, height) * 200000), 50000)
        etree.SubElement(avLst, qn("a:gd"), attrib={"name": "adj", "fmla": f"val {ratio}"})
    return shape


def _apply_gradient(shape, hex1: str, hex2: str, hex3: str, angle_deg: int = 135):
    """3-stop linear gradient fill (OOXML)"""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    for tag in ["a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"]:
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, hx in zip([0, 50000, 100000], [hex1, hex2, hex3]):
        gs = etree.SubElement(gsLst, qn("a:gs"), attrib={"pos": str(pos)})
        etree.SubElement(gs, qn("a:srgbClr"), attrib={"val": hx})

    ang_ooxml = int(angle_deg * 60000) % 21600000
    etree.SubElement(gradFill, qn("a:lin"), attrib={
        "ang": str(ang_ooxml), "scaled": "0"
    })
    shape.line.fill.background()


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_bullets(slide, left, top, width, height, items: list, size=12):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"· {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = C["body"]


# ── 공통 프레임 (모든 슬라이드) ──────────────────────────────────────────────

def _add_frame(slide, chapter: str, page_num: int, dark: bool = False):
    text_c = C["white"] if dark else C["muted"]

    # 챕터명 (좌측)
    _add_textbox(slide, MH, HEADER_TOP, Inches(6), HEADER_H,
                 text=chapter, size=11, color=text_c)

    # 로고 자리 (우측) - "Brandlogy" 텍스트로 표시
    _add_textbox(slide, W - MH - Inches(1.5), HEADER_TOP, Inches(1.5), HEADER_H,
                 text="Brandlogy", size=11, bold=True, color=text_c, align=PP_ALIGN.RIGHT)

    # 페이지 번호 (하단 좌)
    _add_textbox(slide, MH, FOOTER_TOP, Inches(2), FOOTER_H,
                 text=str(page_num), size=10, color=text_c)


def _add_headline_block(slide, headline: str, subtitle: str = ""):
    _add_textbox(slide, MH, HEADLINE_TOP, CONTENT_W, HEADLINE_H,
                 text=headline, size=34, bold=True, color=C["headline"])
    if subtitle:
        _add_textbox(slide, MH, SUBTITLE_TOP, CONTENT_W, SUBTITLE_H,
                     text=subtitle, size=16, color=C["subtitle"])


# ── 슬라이드 빌더 ─────────────────────────────────────────────────────────────

def _cover_slide(prs: Presentation, topic: str, date_str: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, C["white"])

    # 헤더
    _add_textbox(slide, MH, HEADER_TOP, Inches(6), HEADER_H,
                 text="Research Report", size=11, color=C["muted"])
    _add_textbox(slide, W - MH - Inches(1.5), HEADER_TOP, Inches(1.5), HEADER_H,
                 text="Brandlogy", size=11, bold=True, color=C["muted"], align=PP_ALIGN.RIGHT)

    # 타이틀
    _add_textbox(slide, MH, HEADLINE_TOP, CONTENT_W, Inches(1.2),
                 text=topic, size=40, bold=True, color=C["headline"])
    _add_textbox(slide, MH, Inches(2.3), CONTENT_W, Inches(0.4),
                 text="AI 기반 종합 리서치 보고서", size=16, color=C["subtitle"])

    # 히어로 카드 (gradient)
    card = _add_rect(slide,
                     MH, Inches(3.0), CONTENT_W, Inches(3.5),
                     radius_emu=int(Inches(0.22)))
    _apply_gradient(card, "1456f0", "3b82f6", "60a5fa", 135)

    # 카드 내부 텍스트
    _add_textbox(slide, MH + Inches(0.4), Inches(3.6), Inches(8), Inches(0.8),
                 text=topic, size=28, bold=True, color=C["white"])
    _add_textbox(slide, MH + Inches(0.4), Inches(4.5), Inches(8), Inches(0.4),
                 text="종합 리서치 보고서  |  " + date_str, size=13, color=RGBColor(0xbf, 0xdb, 0xfe))

    # 푸터
    _add_textbox(slide, MH, FOOTER_TOP, Inches(3), FOOTER_H,
                 text="1", size=10, color=C["muted"])


def _section_divider(prs: Presentation, section_num: int,
                     section_title: str, page_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C["dark"])

    # 섹션 번호
    _add_textbox(slide, MH, HEADER_TOP, Inches(4), HEADER_H,
                 text=f"Section {section_num:02d}", size=12, color=RGBColor(0xff, 0xff, 0xff))
    _add_textbox(slide, W - MH - Inches(1.5), HEADER_TOP, Inches(1.5), HEADER_H,
                 text="Brandlogy", size=11, bold=True,
                 color=C["white"], align=PP_ALIGN.RIGHT)

    # 섹션 타이틀 (세로 중앙)
    _add_textbox(slide, MH, Inches(2.8), CONTENT_W, Inches(1.4),
                 text=section_title, size=44, bold=True, color=C["white"])

    # 강조 라인
    line = _add_rect(slide, MH, Inches(4.3), Inches(1.5), Inches(0.05),
                     fill_color=C["blue2"])

    _add_textbox(slide, MH, FOOTER_TOP, Inches(2), FOOTER_H,
                 text=str(page_num), size=10, color=RGBColor(0xff, 0xff, 0xff))


def _content_slide(prs: Presentation, chapter: str, headline: str,
                   subtitle: str, bullets: list, page_num: int,
                   sources: list = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C["white"])
    _add_frame(slide, chapter, page_num)
    _add_headline_block(slide, headline, subtitle)

    # 불릿 카드 (최대 5개)
    bullets = bullets[:5]
    if not bullets:
        return

    card_h = Inches(4.0) // len(bullets)
    card_h = min(card_h, Inches(1.1))

    for i, bullet in enumerate(bullets):
        top = BODY_TOP + i * (card_h + Inches(0.08))
        if top + card_h > Inches(6.85):
            break

        bg = _add_rect(slide, MH, top, CONTENT_W, card_h,
                       fill_color=C["surface"], radius_emu=int(Inches(0.12)))

        # 번호 배지
        badge = _add_rect(slide, MH + Inches(0.18), top + Inches(0.18),
                          Inches(0.35), Inches(0.35),
                          fill_color=C["blue"], radius_emu=int(Inches(0.05)))

        _add_textbox(slide, MH + Inches(0.18), top + Inches(0.12),
                     Inches(0.35), Inches(0.35),
                     text=str(i + 1), size=11, bold=True, color=C["white"],
                     align=PP_ALIGN.CENTER)

        _add_textbox(slide, MH + Inches(0.65), top + Inches(0.1),
                     CONTENT_W - Inches(0.8), card_h - Inches(0.2),
                     text=bullet, size=12, color=C["body"], wrap=True)

    # 출처
    if sources:
        src_text = " | ".join(sources[:3])
        _add_textbox(slide, MH, Inches(6.88), CONTENT_W, Inches(0.15),
                     text=f"출처: {src_text}", size=9, color=C["muted"])


def _kpi_slide(prs: Presentation, chapter: str, headline: str,
               subtitle: str, kpis: list, body_text: str,
               page_num: int) -> None:
    """KPI 타일 + 설명 텍스트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C["white"])
    _add_frame(slide, chapter, page_num)
    _add_headline_block(slide, headline, subtitle)

    kpis = kpis[:4]
    if kpis:
        tile_w = CONTENT_W // len(kpis) - Inches(0.15)
        for i, kpi in enumerate(kpis):
            left = MH + i * (tile_w + Inches(0.15))
            card = _add_rect(slide, left, BODY_TOP, tile_w, Inches(1.6),
                             fill_color=C["white"], line_color=C["border"],
                             line_width=Pt(1), radius_emu=int(Inches(0.12)))
            _add_textbox(slide, left + Inches(0.2), BODY_TOP + Inches(0.2),
                         tile_w - Inches(0.4), Inches(0.7),
                         text=kpi.get("value", ""), size=28, bold=True, color=C["blue"])
            _add_textbox(slide, left + Inches(0.2), BODY_TOP + Inches(0.95),
                         tile_w - Inches(0.4), Inches(0.4),
                         text=kpi.get("label", ""), size=11, color=C["subtitle"])

    # 본문 텍스트 (KPI 아래)
    if body_text:
        _add_textbox(slide, MH, BODY_TOP + Inches(1.85), CONTENT_W, Inches(2.5),
                     text=body_text, size=13, color=C["body"], wrap=True)


def _resources_slide(prs: Presentation, chapter: str, page_num: int,
                     papers: list, hbr: list, courses: list) -> None:
    """학습 자료 큐레이션 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C["white"])
    _add_frame(slide, chapter, page_num)
    _add_headline_block(slide, "학습 자료 큐레이션",
                        "핵심 논문 · HBR 아티클 · 온라인 강의")

    col_w = Inches(3.8)
    cols = [
        ("📄 핵심 논문", papers[:3],
         lambda p: f"{p.get('year','')} {p.get('title','')[:40]}"),
        ("📰 HBR 아티클", hbr[:4],
         lambda a: a.get("title","")[:45]),
        ("🎓 온라인 강의", courses[:3],
         lambda c: f"[{c.get('platform','')}] {c.get('title','')[:35]}"),
    ]

    for ci, (col_title, items, fmt) in enumerate(cols):
        left = MH + ci * (col_w + Inches(0.22))

        # 컬럼 헤더
        header = _add_rect(slide, left, BODY_TOP, col_w, Inches(0.4),
                           fill_color=C["blue"], radius_emu=int(Inches(0.1)))
        _add_textbox(slide, left + Inches(0.15), BODY_TOP + Inches(0.05),
                     col_w - Inches(0.3), Inches(0.35),
                     text=col_title, size=12, bold=True, color=C["white"])

        for ri, item in enumerate(items):
            item_top = BODY_TOP + Inches(0.55) + ri * Inches(1.05)
            if item_top + Inches(0.9) > Inches(6.85):
                break
            card = _add_rect(slide, left, item_top, col_w, Inches(0.95),
                             fill_color=C["surface"], radius_emu=int(Inches(0.1)))
            _add_textbox(slide, left + Inches(0.15), item_top + Inches(0.1),
                         col_w - Inches(0.3), Inches(0.75),
                         text=fmt(item), size=10, color=C["body"], wrap=True)


def _closing_slide(prs: Presentation, topic: str, page_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C["dark"])

    _add_textbox(slide, W - MH - Inches(1.5), HEADER_TOP, Inches(1.5), HEADER_H,
                 text="Brandlogy", size=11, bold=True,
                 color=C["white"], align=PP_ALIGN.RIGHT)

    _add_textbox(slide, MH, Inches(2.5), CONTENT_W, Inches(1.0),
                 text="Thank You", size=48, bold=True, color=C["white"],
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, MH, Inches(3.7), CONTENT_W, Inches(0.5),
                 text=topic, size=18, color=RGBColor(0x93, 0xc5, 0xfd),
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, MH, FOOTER_TOP, Inches(2), FOOTER_H,
                 text=str(page_num), size=10, color=RGBColor(0xff, 0xff, 0xff))


# ── 마크다운 파서 ─────────────────────────────────────────────────────────────

def _parse_content(md_text: str) -> dict:
    lines = md_text.split("\n")
    bullets, h2s, paragraphs = [], [], []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("## "):
            h2s.append(stripped[3:])
        elif stripped.startswith("### "):
            h2s.append(stripped[4:])
        elif stripped.startswith(("- ", "* ", "· ")):
            bullets.append(stripped[2:].strip())
        elif not stripped.startswith("#"):
            paragraphs.append(stripped)

    # 첫 문장 → subtitle
    subtitle = ""
    for p in paragraphs:
        if len(p) > 20:
            subtitle = p[:80] + ("..." if len(p) > 80 else "")
            break

    return {
        "bullets": bullets[:8],
        "h2s": h2s[:6],
        "paragraphs": paragraphs[:4],
        "subtitle": subtitle,
    }


def _extract_kpis(md_text: str) -> list:
    """숫자 + 단위 패턴을 KPI로 추출"""
    pattern = r'(\d+[\d,\.]*\s*(?:%|명|개|편|개월|년|억|배|점|위))'
    matches = re.findall(pattern, md_text)
    kpis = []
    for m in matches[:4]:
        idx = md_text.find(m)
        context = md_text[max(0, idx-20):idx].strip()
        label = re.sub(r'[^가-힣a-zA-Z\s]', '', context)[-15:].strip() or "Key Metric"
        kpis.append({"value": m.strip(), "label": label})
    return kpis


# ── 메인 ─────────────────────────────────────────────────────────────────────

SECTION_INFO = {
    "의미":    ("1", "의미와 중요성"),
    "방법":    ("2", "실행 방법론"),
    "사례":    ("3", "AI 시대 대표 사례"),
    "학습자료": ("4", "학습 자료 큐레이션"),
}


def run(synthesis: dict) -> str:
    topic = synthesis.get("topic", "리서치")
    sections = synthesis.get("sections", {})
    sources = synthesis.get("sources", {})
    papers = sources.get("papers", [])
    hbr = sources.get("hbr_articles", [])
    courses = sources.get("courses", [])

    print(f"[PPT 에이전트] '{topic}' 프레젠테이션 생성 시작...")
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 빈 레이아웃 확보
    while len(prs.slide_layouts) < 7:
        prs.slide_layouts.add_slide_layout(prs.slide_layouts[0])

    page = 1
    date_str = datetime.now().strftime("%Y.%m.%d")

    # 1. 표지
    _cover_slide(prs, topic, date_str)
    page += 1

    # 2. 섹션별 슬라이드
    section_order = ["의미", "방법", "사례", "학습자료"]
    for key in section_order:
        if key not in sections:
            continue

        sec = sections[key]
        sec_num, sec_label = SECTION_INFO.get(key, ("?", key))
        content = sec.get("content", "")
        title = sec.get("title", sec_label)
        parsed = _parse_content(content)
        kpis = _extract_kpis(content)

        # 섹션 디바이더
        _section_divider(prs, int(sec_num), title, page)
        page += 1

        if key == "학습자료":
            _resources_slide(prs, title, page, papers, hbr, courses)
            page += 1
        else:
            # KPI 슬라이드 (숫자 데이터 있을 때)
            if kpis:
                _kpi_slide(prs, title, title, parsed["subtitle"],
                           kpis, "\n".join(parsed["paragraphs"][:2]), page)
                page += 1

            # 본문 불릿 슬라이드
            bullets = parsed["bullets"] if parsed["bullets"] else parsed["h2s"][:5]
            if bullets:
                _content_slide(prs, title, title, parsed["subtitle"],
                               bullets, page)
                page += 1

    # 클로징
    _closing_slide(prs, topic, page)

    # 저장
    today = datetime.now().strftime("%Y%m%d_%H%M")
    safe = re.sub(r'[\\/*?:"<>|]', "_", topic)[:30]
    out_path = os.path.join(OUTPUT_DIR, f"{safe}_{today}.pptx")
    prs.save(out_path)

    print(f"[PPT 에이전트] 저장 완료 → {out_path}")
    return out_path


if __name__ == "__main__":
    import json, sys
    data_dir = os.environ.get("RESEARCH_DATA_DIR", os.path.join(BASE_DIR, "data"))
    syn_path = os.path.join(data_dir, "synthesis.json")
    if not os.path.exists(syn_path):
        print("synthesis.json 없음.")
        sys.exit(1)
    with open(syn_path, encoding="utf-8") as f:
        synthesis = json.load(f)
    run(synthesis)
